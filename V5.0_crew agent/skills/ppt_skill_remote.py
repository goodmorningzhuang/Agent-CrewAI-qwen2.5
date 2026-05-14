import os
import re
import requests
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from datetime import datetime


def call_api_refine(text_content, note, api_base, api_key, model_name):
    """
    使用外部 API 模型进行硬件专家级 PPT 文案润色
    一次性处理所有文字框内容，提高效率
    """
    prompt = f"""你是一位资深的硬件工程专家，同时也是PPT文案设计大师。
请根据以下【特定备注要求】对技术内容进行专家级润色和完善。

【特定备注要求】：{note}

严格规则：
1. 使用硬件工程领域的专业术语（如：性能稳定性、兼容性验证、信号完整性、EMC/EMI、热设计功耗等）。
2. 对原文进行详尽优化、补充和完善，但必须压缩为简洁的要点形式。
3. 保持 Bullet Points 格式，每条要点不超过 25 个字。
4. 如果原文信息不足，请基于硬件专业知识进行合理补充，使内容更加完整和专业。
5. 严禁使用口语化表达，必须使用工业级书面语言。
6. 直接输出结果，不要任何开场白、解释或标注。

待润色内容：
{text_content}

专业润色后："""

    url = f"{api_base}/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": model_name,
        "messages": [{"role": "user", "content": prompt}],
        "stream": False,
        "max_tokens": 512
    }
    try:
        response = requests.post(url, json=payload, headers=headers, timeout=60)
        if response.status_code == 200:
            result = response.json()["choices"][0]["message"]["content"].strip()
            print(f"  {GREEN}✔ API 润色成功 (输入{len(text_content)}字 → 输出{len(result)}字){RESET}")
            return result
        else:
            print(f"  ⚠ API 返回状态码 {response.status_code}，保留原文")
            return text_content
    except Exception as e:
        print(f"  ⚠ API 调用异常: {e}，保留原文")
        return text_content


# 颜色常量（与 main.py 保持一致）
GREEN = '\033[92m'
YELLOW = '\033[93m'
CYAN = '\033[96m'
RESET = '\033[0m'


def _calc_font_size(text, pref_size, shape_width, shape_height):
    """
    根据文字长度和文本框尺寸智能计算字号
    目标：文字能完整显示在文本框内
    """
    char_count = len(text)
    
    # 基础字号：根据文字长度调整
    if char_count > 100:
        base_size = max(pref_size - 5, 7)
    elif char_count > 80:
        base_size = max(pref_size - 4, 8)
    elif char_count > 60:
        base_size = max(pref_size - 3, 8)
    elif char_count > 40:
        base_size = max(pref_size - 2, 9)
    elif char_count > 25:
        base_size = pref_size
    elif char_count > 15:
        base_size = pref_size + 1
    else:
        base_size = pref_size + 3
    
    # 根据文本框面积进一步调整
    if shape_width and shape_height:
        # 计算文本框的大致可用面积（EMU单位）
        area = shape_width * shape_height
        # 估算每行可容纳的字符数（假设字体宽高比约0.6）
        char_width = base_size * 0.6 * 12700  # Pt to EMU 近似
        if char_width > 0:
            chars_per_line = max(int(shape_width / char_width), 1)
        else:
            chars_per_line = 20
        
        # 估算行数
        estimated_lines = max(1, char_count // chars_per_line + 1)
        line_height = base_size * 1.5 * 12700  # 行间距
        total_height_needed = estimated_lines * line_height
        
        # 如果需要的高度超过文本框高度，缩小字号
        if total_height_needed > shape_height and shape_height > 0:
            # 计算缩放比例
            ratio = shape_height / total_height_needed
            adjusted_size = max(int(base_size * ratio), 7)
            base_size = adjusted_size
        
        # 如果文字很少但文本框很小，也要适当缩小
        if area < 2000000000 and char_count > 30:  # 约 5cm x 5cm 的区域
            base_size = max(base_size - 1, 7)
    
    return base_size


def _write_text_to_shape(shape, text, pref_size):
    """将润色后的文字写回文本框，智能匹配框大小"""
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    
    # 清空原有内容
    for i in range(len(text_frame.paragraphs)):
        p = text_frame.paragraphs[i]
        p.clear()
    
    # 记录原始文本框尺寸
    shape_width = shape.width
    shape_height = shape.height
    
    # 按行分割文字
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    if not lines:
        lines = [text]
    
    # 去掉markdown格式的 * 号
    cleaned_lines = []
    for line in lines:
        line = line.lstrip(' ')
        line = line.replace('**', '').replace('*', '')
        if line.startswith('- ') or line.startswith('• '):
            line = line[2:]
        cleaned_lines.append(line)
    lines = cleaned_lines
    
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        p.font.name = "微软雅黑"
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        
        # 计算智能字号
        font_size = _calc_font_size(line, pref_size, shape_width, shape_height)
        p.font.size = Pt(font_size)


def run_ppt_skill(file_path, *args, **kwargs):
    """
    AI增强版 PPT 处理：
    1. 遍历每一页的每个文字框，提取文字
    2. 汇总所有文字一次性发给外部API进行硬件专家级润色
    3. 将润色后的文字根据文字框大小智能缩放后回写
    """
    target = os.path.abspath(file_path)
    if not os.path.exists(target):
        return f"❌ 找不到文件：{target}"

    # 提取参数
    pref_font_size = int(kwargs.get('font_size', 12))
    file_note = kwargs.get('note', '无特定要求')
    api_base = kwargs.get('api_base', '')
    api_key = kwargs.get('api_key', '')
    model_name = kwargs.get('model_name', '')

    if not api_base or not api_key or not model_name:
        return "❌ 缺少外部 API 配置，无法使用 AI 增强模式"

    try:
        print(f"⚡ 启动 AI 增强版 PPT 优化引擎...")
        print(f"📖 调取档案记忆：字号基准={pref_font_size}pt, 备注={file_note}")
        print(f"🤖 使用模型: {model_name}")

        prs = Presentation(target)
        slide_count = len(prs.slides)

        # ================= 第一步：提取所有文字框的文字 =================
        # 结构: [(slide_index, shape_index, original_text), ...]
        text_boxes = []
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    original_text = shape.text.strip()
                    if len(original_text) > 0:  # 只要有文字就提取
                        text_boxes.append((slide_idx, shape_idx, original_text))

        if not text_boxes:
            return "⚠️ PPT 中未找到任何文字内容"

        print(f"📋 共提取 {len(text_boxes)} 个文字框，开始 AI 润色...")

        # ================= 第二步：汇总文字发给 API 润色 =================
        # 构造一个带编号的汇总文本
        combined_input = ""
        for idx, (s_idx, sh_idx, text) in enumerate(text_boxes):
            combined_input += f"[{idx + 1}] 第{s_idx + 1}页:\n{text}\n\n"

        combined_prompt = f"""你是一位资深的硬件工程专家，同时也是PPT文案设计大师。
请根据以下【特定备注要求】对所有文字框内容进行专家级润色和完善。

【特定备注要求】：{file_note}

严格规则：
1. 使用硬件工程领域的专业术语（如：性能稳定性、兼容性验证、信号完整性、EMC/EMI、热设计功耗等）。
2. 对每条内容进行详尽优化、补充和完善，但必须压缩为简洁的要点形式。
3. 保持每条要点简洁，每条不超过 25 个字。
4. 如果原文信息不足，请基于硬件专业知识进行合理补充。
5. 严禁使用口语化表达，必须使用工业级书面语言。
6. **必须保持编号格式**：每条结果以 [编号] 开头，与输入编号一一对应。
7. 如果某个文字框是标题（很短），润色后也要保持简短。
8. 直接输出结果，不要任何开场白、解释或标注。

待润色内容：
{combined_input}

请逐条润色，保持编号对应："""

        print(f"🤖 正在调用外部 API 进行硬件专家级润色...")
        url = f"{api_base}/chat/completions"
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": model_name,
            "messages": [{"role": "user", "content": combined_prompt}],
            "stream": False,
            "max_tokens": 2048
        }

        try:
            response = requests.post(url, json=payload, headers=headers, timeout=120)
            if response.status_code == 200:
                combined_result = response.json()["choices"][0]["message"]["content"].strip()
                print(f"{GREEN}✔ API 润色完成 (输入{len(combined_input)}字 → 输出{len(combined_result)}字){RESET}")
            else:
                print(f"{YELLOW}⚠ API 返回状态码 {response.status_code}，保留原文{RESET}")
                combined_result = None
        except Exception as e:
            print(f"{YELLOW}⚠ API 调用异常: {e}，保留原文{RESET}")
            combined_result = None

        # ================= 第三步：解析润色结果 =================
        if combined_result:
            # 按编号解析: [1] xxx  [2] xxx  ...
            # 匹配 [数字] 开头的内容
            parts = re.split(r'\[(\d+)\]', combined_result)
            # parts: ['', '1', '内容1', '2', '内容2', ...]
            refined_texts = {}
            for i in range(1, len(parts) - 1, 2):
                idx = int(parts[i])
                content = parts[i + 1].strip()
                if content:
                    refined_texts[idx] = content
            
            print(f"📝 成功解析 {len(refined_texts)} 条润色结果")
        else:
            refined_texts = {}

        # ================= 第四步：将润色后的文字写回对应的文字框 =================
        update_count = 0
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if not (hasattr(shape, "text_frame") and shape.text_frame is not None):
                    continue
                original_text = shape.text.strip()
                if not original_text:
                    continue

                # 找到对应的编号
                for box_idx, (s_idx, sh_idx, orig_text) in enumerate(text_boxes):
                    if s_idx == slide_idx and sh_idx == shape_idx and orig_text == original_text:
                        box_number = box_idx + 1
                        if box_number in refined_texts:
                            refined_text = refined_texts[box_number]
                            print(f"  📝 第{slide_idx + 1}页 文字框{box_number}: \"{orig_text[:20]}...\" → \"{refined_text[:20]}...\"")
                            _write_text_to_shape(shape, refined_text, pref_font_size)
                            update_count += 1
                        break

        print(f"\n📊 共更新 {update_count}/{len(text_boxes)} 个文字框")

        # ================= 保存文件 =================
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        file_dir, full_name = os.path.split(target)
        name_part, ext_part = os.path.splitext(full_name)
        new_filename = f"{name_part}_AI_Enhanced_{timestamp}{ext_part}"
        save_path = os.path.join(file_dir, new_filename)

        prs.save(save_path)
        return f"✨ AI增强版 PPT 优化完成！\n📂 备注应用：{file_note}\n📊 更新了 {update_count}/{len(text_boxes)} 个文字框\n💾 文件已生成：{os.path.basename(save_path)}"

    except Exception as e:
        import traceback
        traceback.print_exc()
        return f"❌ AI 增强处理失败: {str(e)}"