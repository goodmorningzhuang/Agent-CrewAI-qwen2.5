import os
import requests
import time
from pptx import Presentation
from datetime import datetime

# 路径指向 PPT 专用文件夹
PPT_DIR = r"C:\Users\user\Desktop\test-ai\PPT"
OLLAMA_URL = "http://localhost:11434/api/generate"
MODEL_NAME = "qwen2.5:1.5b"

def ask_qwen(text):
    if not text.strip() or len(text) < 2: return text
    prompt = f"把这段PPT文字改写得更专业商务，直接给出结果：\n{text}"
    payload = {"model": MODEL_NAME, "prompt": prompt, "stream": False}
    try:
        response = requests.post(OLLAMA_URL, json=payload, timeout=60)
        return response.json().get('response', '').strip().replace('"', '')
    except:
        return text

def run_ppt_skill(filename):
    input_path = os.path.join(PPT_DIR, filename)
    today_str = datetime.now().strftime("%Y%m%d")
    output_path = os.path.join(PPT_DIR, f"优化完成_{today_str}_{filename}")

    if not os.path.exists(input_path):
        print(f"❌ 找不到文件: {input_path}")
        return

    print(f"🚀 正在润色...")
    prs = Presentation(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for paragraph in shape.text_frame.paragraphs:
                if len(paragraph.text.strip()) > 2:
                    # 字体保留逻辑
                    old_size = None
                    for run in paragraph.runs:
                        if run.font.size: 
                            old_size = run.font.size
                            break
                    
                    paragraph.text = ask_qwen(paragraph.text)
                    
                    if old_size:
                        for run in paragraph.runs: run.font.size = old_size
    
    prs.save(output_path)
    print(f"✅ 完成！保存位置: {output_path}")