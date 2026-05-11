import os
import requests
from pptx import Presentation
from pptx.util import Pt
from datetime import datetime

# 本地模型配置
OLLAMA_URL = "http://localhost:11434/api/generate"
MODEL_NAME = "qwen2.5:1.5b"

def call_qwen_for_refined_ppt(text, note="无"):
    """
    专家级精炼：结合【文件备注】进行极致压缩
    """
    # 将长记忆备注注入提示词
    prompt = f"""
    你是一个硬件工程领域的 PPT 文案专家。
    请根据以下【特定备注要求】对技术内容进行专家级润色：
    
    【特定备注要求】：{note}
    
    规则：
    1. 必须使用 Bullet Points 形式。
    2. 总字数控制在 40 字以内。
    3. 使用专业词汇（如：性能稳定性、兼容性验证）。
    4. 严禁长句，直接输出结果。

    待处理内容：
    {text}
    
    精炼后：
    """
    try:
        response = requests.post(OLLAMA_URL, json={
            "model": MODEL_NAME, "prompt": prompt, "stream": False
        }, timeout=30)
        return response.json().get("response", "").strip()
    except:
        return text

def run_ppt_skill(file_path, *args, **kwargs):
    """
    适配 Agent 3.0：接收全局字号偏好和文件级长记忆
    """
    target = os.path.abspath(file_path)
    if not os.path.exists(target):
        return f"❌ 找不到文件：{target}"

    # 1. 提取长记忆参数
    # 获取全局偏好的字号，默认为 12
    pref_font_size = int(kwargs.get('font_size', 12))
    # 获取特定文件的备注要求
    file_note = kwargs.get('note', '无特定要求')

    try:
        print(f"⚡ 启动全页自适应优化引擎...")
        print(f"📖 调取档案记忆：字号基准={pref_font_size}pt, 备注={file_note}")
        
        prs = Presentation(target)
        
        for index, slide in enumerate(prs.slides):
            print(f"📝 正在同步优化第 {index+1} 页...", end="\r")
            for shape in slide.shapes:
                if hasattr(shape, "text") and len(shape.text.strip()) > 15:
                    original_text = shape.text
                    # 传入长记忆备注进行文案润色
                    refined_text = call_qwen_for_refined_ppt(original_text, note=file_note)
                    
                    shape.text = "" 
                    text_frame = shape.text_frame
                    text_frame.word_wrap = True 
                    
                    p = text_frame.paragraphs[0]
                    p.text = refined_text
                    p.font.name = "微软雅黑"
                    
                    # --- 核心：基于全局偏好的字号自适应算法 ---
                    char_len = len(refined_text)
                    if char_len > 100:
                        p.font.size = Pt(pref_font_size - 2) # 文字极多时比偏好小 2 号
                    elif char_len > 50:
                        p.font.size = Pt(pref_font_size)     # 匹配偏好字号
                    else:
                        p.font.size = Pt(pref_font_size + 2) # 文字精简时比偏好大 2 号

        # 生成带时间戳的文件名
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        file_dir, full_name = os.path.split(target)
        name_part, ext_part = os.path.splitext(full_name)
        new_filename = f"{name_part}_AutoFix_{timestamp}{ext_part}"
        save_path = os.path.join(file_dir, new_filename)
        
        prs.save(save_path)
        return f"✨ 全页优化成功！\n📂 备注应用：{file_note}\n💾 文件已生成：{os.path.basename(save_path)}"

    except Exception as e:
        return f"❌ 处理失败: {str(e)}"