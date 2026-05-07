import os
import requests
import time
from pptx import Presentation
from datetime import datetime

# --- 环境配置 ---
BASE_PATH = r"C:\Users\user\Desktop\test-ai"
INPUT_FILENAME = "AI硬件BU测试0409.pptx" 
INPUT_FILE = os.path.join(BASE_PATH, INPUT_FILENAME)

# 输出文件名加上当前日期
today_str = datetime.now().strftime("%Y%m%d")
OUTPUT_FILENAME = f"优化完成_AI硬件BU测试_{today_str}.pptx"
OUTPUT_FILE = os.path.join(BASE_PATH, OUTPUT_FILENAME)

OLLAMA_URL = "http://localhost:11434/api/generate"
# 修改为更轻量的 1.5b 模型，速度更快
MODEL_NAME = "qwen2.5:1.5b" 

def ask_qwen(text):
    """利用 Qwen 2.5:1.5b 进行快速文案优化"""
    if not text.strip() or len(text) < 2:
        return text
    
    # 针对 1.5b 模型的 Prompt：指令更加直接
    prompt = (
        f"把这段PPT文字改写得更专业、更商务，直接给出结果，不要解释。\n"
        f"原文：{text}\n"
        f"改写："
    )
    
    payload = {
        "model": MODEL_NAME,
        "prompt": prompt,
        "stream": False,
        "options": {
            "temperature": 0.5,
            "num_predict": 100 # 限制长度以确保速度
        }
    }
    
    try:
        # 1.5b 模型通常很快，timeout 设为 60s 绰绰有余
        response = requests.post(OLLAMA_URL, json=payload, timeout=60)
        improved = response.json().get('response', '').strip()
        # 清洗可能存在的引号
        return improved.replace('"', '').replace('“', '').replace('”', '')
    except Exception as e:
        print(f"⚠️ 网络请求异常: {e}")
        return text

def run_agent():
    if not os.path.exists(INPUT_FILE):
        print(f"❌ 错误：在 {BASE_PATH} 没找到 {INPUT_FILENAME}")
        return

    start_time = time.time()
    print(f"🚀 启动 {MODEL_NAME} 极速版 Agent...")
    
    try:
        prs = Presentation(INPUT_FILE)
    except Exception as e:
        print(f"❌ 无法打开文件（请确认 PPT 未被占用）：{e}")
        return
    
    slide_count = len(prs.slides)
    modified_count = 0

    try:
        for i, slide in enumerate(prs.slides):
            print(f"正在处理第 {i+1}/{slide_count} 页...")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    original_text = paragraph.text
                    if len(original_text.strip()) > 2:
                        improved_text = ask_qwen(original_text)

                        # 获取原始第一个 run 的字体大小作为参考
                        original_font_size = None
                        for run in paragraph.runs:
                            if run.font.size is not None:
                                original_font_size = run.font.size
                                break
                        
                        if improved_text and improved_text != original_text:
                            print(f"  ✅ 修改: {original_text[:10].strip()}... -> {improved_text[:10].strip()}...")
                            paragraph.text = improved_text
                            modified_count += 1
                            
                            # 恢复原始字体大小
                            if original_font_size is not None:
                                for run in paragraph.runs:
                                    run.font.size = original_font_size
        
        # 尝试保存文件
        prs.save(OUTPUT_FILE)
        duration = time.time() - start_time
        print("\n" + "="*30)
        print(f"✨ 极速优化完成！")
        print(f"⏱️ 总耗时: {duration:.2f} 秒")
        print(f"📝 修改位置: {modified_count} 处")
        print(f"💾 保存路径: {OUTPUT_FILE}")
        print("="*30)

    except Exception as e:
        print(f"❌ 运行过程中出现错误: {e}")

if __name__ == "__main__":
    run_agent()