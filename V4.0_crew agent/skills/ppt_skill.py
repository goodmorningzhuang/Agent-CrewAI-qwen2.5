import os
import requests
from pptx import Presentation
from pptx.util import Pt
from datetime import datetime

# 本地模型配置
OLLAMA_URL = "http://localhost:11434/api/generate"
MODEL_NAME = "qwen2.5:1.5b"

def call_qwen_for_refined_ppt(text):
    """
    专家级精炼：保持硬件专业性的同时极致压缩
    """
    prompt = f"""
    你是一个硬件工程领域的 PPT 文案专家。请将以下技术内容进行【专家级润色】并【极致精炼】：
    1. 必须使用 Bullet Points 形式。
    2. 总字数控制在 40 字以内。
    3. 使用专业词汇（如：性能稳定性、兼容性验证、参数调优）。
    4. 严禁长句，直接输出结果，不要解释。

    待处理内容：
    {text}
    
    精炼后：
    """
    try:
        response = requests.post(OLLAMA_URL, json={
            "model": MODEL_NAME, "prompt": prompt, "stream": False
        }, timeout=30)
        return response.json().get("response", "").strip()
    except:
        return text

def run_ppt_skill(file_path):
    target = os.path.abspath(file_path)
    if not os.path.exists(target):
        return f"❌ 找不到文件：{target}"

    try:
        print(f"⚡ 启动全页自适应优化引擎...")
        prs = Presentation(target)
        
        # 遍历每一页（解决“后面没处理”的问题）
        for index, slide in enumerate(prs.slides):
            print(f"📝 正在同步优化第 {index+1} 页...")
            for shape in slide.shapes:
                if hasattr(shape, "text") and len(shape.text.strip()) > 15:
                    original_text = shape.text
                    refined_text = call_qwen_for_refined_ppt(original_text)
                    
                    # 彻底清空，防止文字叠加
                    shape.text = "" 
                    text_frame = shape.text_frame
                    text_frame.word_wrap = True # 开启自动换行
                    
                    p = text_frame.paragraphs[0]
                    p.text = refined_text
                    p.font.name = "微软雅黑"
                    
                    # --- 核心：自动缩小文字匹配原本的框 ---
                    char_len = len(refined_text)
                    if char_len > 100:
                        p.font.size = Pt(10) # 极多文字时用小号
                    elif char_len > 50:
                        p.font.size = Pt(12) # 较多文字
                    else:
                        p.font.size = Pt(14) # 精简文字

        # 生成带时间戳的文件名
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        name_part, ext_part = os.path.splitext(target)
        new_filename = f"{name_part}_AutoFix_{timestamp}{ext_part}"
        
        prs.save(new_filename)
        return f"✨ 全页优化成功！\n💾 文件已生成：{os.path.basename(new_filename)}"

    except Exception as e:
        return f"❌ 处理失败: {str(e)}"