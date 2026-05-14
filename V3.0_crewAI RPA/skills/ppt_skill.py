import os
import requests
from pptx import Presentation
from datetime import datetime
from pptx.util import Pt  # 精准控制字号

def run_ppt_skill(file_path):
    # 使用绝对路径，彻底根治 ppt/ppt 叠加报错
    target = os.path.abspath(file_path)
    if not os.path.exists(target):
        return f"❌ 找不到文件：{target}"

    try:
        prs = Presentation(target)
        total_slides = len(prs.slides)
        print(f"📂 已加载: {os.path.basename(target)} (共 {total_slides} 页)")

        for index, slide in enumerate(prs.slides, start=1):
            print(f"   进度: [{index}/{total_slides}] - 正在进行专家级润色与排版自适应...", end='\r')
            
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if shape.text.strip():
                        # 备份原始字体与颜色
                        try:
                            first_run = shape.text_frame.paragraphs[0].runs[0]
                            orig_font_name = first_run.font.name
                            orig_font_color = first_run.font.color.rgb if first_run.font.color else None
                        except:
                            orig_font_name = "微软雅黑"
                            orig_font_color = None

                        # 高质量润色 Prompt
                        prompt = f"""
                        你是一个资深的硬件自动化专家。请对以下文案进行专业润色。
                        要求：
                        - 使用严谨的硬件术语（如：NPI、量产稳定性、自动化覆盖率）。
                        - 提升逻辑性，使其符合专业汇报水平。
                        - 润色后内容需保持精炼，严禁废话。
                        
                        待润色文案：{shape.text}
                        
                        直接输出润色结果：
                        """

                        try:
                            response = requests.post(
                                "http://localhost:11434/api/generate",
                                json={"model": "qwen2.5:1.5b", "prompt": prompt, "stream": False},
                                timeout=20
                            )
                            refined_text = response.json().get("response", shape.text).strip()
                            
                            # 核心：物理适配逻辑
                            tf = shape.text_frame
                            tf.word_wrap = True
                            tf.text = refined_text
                            
                            # 根据字数强行分配字号，防止重叠
                            text_len = len(refined_text)
                            if text_len > 80:
                                target_size = Pt(9)   # 极长文案使用超小号
                            elif text_len > 50:
                                target_size = Pt(10.5)
                            elif text_len > 30:
                                target_size = Pt(12)
                            else:
                                target_size = Pt(14)  # 短文案保持大号

                            for paragraph in tf.paragraphs:
                                for run in paragraph.runs:
                                    run.font.name = orig_font_name
                                    run.font.size = target_size
                                    if orig_font_color:
                                        run.font.color.rgb = orig_font_color
                        except:
                            continue
        
        # 带时间戳保存
        time_suffix = datetime.now().strftime("%Y%m%d_%H%M%S")
        file_dir, full_name = os.path.split(target)
        name_part, ext_part = os.path.splitext(full_name)
        new_filename = f"{name_part}_专家适配版_{time_suffix}{ext_part}"
        save_path = os.path.join(file_dir, new_filename)

        prs.save(save_path)
        return f"✨ 处理成功！\n💾 适配版已保存: {new_filename}"

    except Exception as e:
        return f"\n❌ 错误: {str(e)}"